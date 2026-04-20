from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form, Request
from sqlalchemy.orm import Session
import uuid
import traceback  # ספרייה קריטית לשליפת מקור השגיאה המדויק
import io
import PyPDF2
from app.core.database import get_db
from app.models import schemas, db_models
from app.agents.graph import exam_generation_app
import docx
from pptx import Presentation
router = APIRouter()


@router.post("/upload")
async def upload_material(
        request: Request,
        text_content: str = Form(None),
        file: UploadFile = File(None),
        db: Session = Depends(get_db)
):
    print("\n" + "=" * 40)
    print("--- 📥 REQUEST RECEIVED: /upload ---")

    try:
        content = ""

        # בדיקה האם המשתמש העלה קובץ
        if file and file.filename:
            print(f"  -> Processing file: {file.filename}")
            file_extension = file.filename.split('.')[-1].lower()

            # קריאת הקובץ פעם אחת לזיכרון (חוסך כפילויות בקוד)
            file_bytes = await file.read()

            # קריאת קובץ PDF
            if file_extension == 'pdf':
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
                for page in pdf_reader.pages:
                    extracted_text = page.extract_text()
                    if extracted_text:
                        content += extracted_text + "\n"
                print(f"  -> ✅ PDF extracted successfully. ({len(pdf_reader.pages)} pages)")

            # קריאת קובץ Word (DOCX)
            elif file_extension == 'docx':
                doc = docx.Document(io.BytesIO(file_bytes))
                # שולף את הטקסט רק מפיסקאות שאינן ריקות
                content = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
                print(f"  -> ✅ DOCX extracted successfully. ({len(doc.paragraphs)} paragraphs)")

            # קריאת קובץ PowerPoint (PPTX)
            elif file_extension == 'pptx':
                prs = Presentation(io.BytesIO(file_bytes))
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_runs.append(shape.text)
                content = "\n".join(text_runs)
                print(f"  -> ✅ PPTX extracted successfully. ({len(prs.slides)} slides)")

            # קריאת קובץ טקסט רגיל
            elif file_extension == 'txt':
                content = file_bytes.decode('utf-8', errors='ignore')
                print("  -> ✅ TXT extracted successfully.")

            else:
                raise HTTPException(status_code=400, detail="פורמט הקובץ אינו נתמך (נא להעלות PDF, DOCX, PPTX או TXT)")

        # אם לא הועלה קובץ, נבדוק אם יש טקסט חופשי
        elif text_content and text_content.strip():
            content = text_content
            print("  -> ✅ Text content mapped successfully.")
        else:
            print("  -> ERROR: No text or file provided.")
            raise HTTPException(status_code=400, detail="חובה לספק טקסט חופשי או קובץ")

        # שמירה במסד הנתונים
        session_id = str(uuid.uuid4())
        print(f"  -> Generated session_id: {session_id}")

        new_exam = db_models.DBExam(
            id=session_id,
            title="טיוטת מבחן",
            source_material_summary=content
        )

        db.add(new_exam)
        db.commit()
        print("  -> ✅ DB Commit successful!")
        print("=" * 40 + "\n")

        return {"session_id": session_id, "message": "החומר הועלה ונשמר בהצלחה!"}

    except HTTPException:
        # מאפשר לשגיאות יזומות (כמו פורמט לא נתמך) לעלות ללקוח בצורה נקייה
        raise
    except Exception as e:
        print("\n🔥 --- CRITICAL ERROR IN /upload --- 🔥")
        print(f"Error Message: {str(e)}")
        db.rollback()
        raise HTTPException(status_code=500, detail=f"שגיאה בהעלאת החומר: {str(e)}")
@router.post("/generate-exam", response_model=schemas.Exam)
async def generate_exam(
        request: schemas.ExamCreateRequest,
        db: Session = Depends(get_db)
):
    print("\n" + "=" * 40)
    print(f"--- 🚀 AI WORKFLOW STARTED: {request.session_id} ---")

    try:
        exam_record = db.query(db_models.DBExam).filter(db_models.DBExam.id == request.session_id).first()
        if not exam_record:
            raise HTTPException(status_code=404, detail="Session_id לא נמצא במערכת")

        initial_state = {
            "session_id": request.session_id,
            "source_text": exam_record.source_material_summary,
            "custom_instructions": request.custom_instructions,
            "blueprint": None,
            "questions": [],
            "review_feedback": None,
            "iteration_count": 0
        }

        final_state = await exam_generation_app.ainvoke(initial_state)

        questions_data = final_state.get("questions", [])
        if not questions_data:
            raise HTTPException(status_code=500, detail="הסוכנים לא הצליחו לייצר שאלות.")

        blueprint = final_state.get("blueprint", {})
        topic_name = blueprint.get("topics", [{"topic_name": "מבחן אישי"}])[0].get("topic_name", "מבחן אישי")
        exam_record.title = f"מבחן בנושא: {topic_name}"

        response_questions = []
        for q_data in questions_data:
            q_id = str(uuid.uuid4())
            new_q = db_models.DBQuestion(
                id=q_id,
                exam_id=request.session_id,
                type=q_data["type"],
                question_text=q_data["question_text"],
                options=q_data.get("options", []),
                correct_answer=q_data["correct_answer"],
                explanation=q_data["explanation"],
                difficulty=q_data["difficulty"]
            )
            db.add(new_q)

            response_questions.append(schemas.Question(
                id=q_id,
                type=q_data["type"],
                question_text=q_data["question_text"],
                options=q_data.get("options", []),
                correct_answer=q_data["correct_answer"],
                explanation=q_data["explanation"],
                difficulty=q_data["difficulty"]
            ))

        db.commit()
        print("  -> ✅ AI Generation and DB save complete!")
        print("=" * 40 + "\n")

        return schemas.Exam(
            id=exam_record.id,
            title=exam_record.title,
            questions=response_questions,
            source_material_summary="המבחן נוצר בהצלחה."
        )

    except Exception as e:
        print("\n🔥 --- CRITICAL ERROR IN /generate-exam --- 🔥")
        print(f"Error Message: {str(e)}")
        print(traceback.format_exc())
        print("=" * 40 + "\n")
        db.rollback()
        raise HTTPException(status_code=500, detail="שגיאה פנימית בהפעלת הסוכנים")


@router.get("/exam/{session_id}")
async def get_exam(session_id: str, db: Session = Depends(get_db)):
    print("\n" + "=" * 40)
    print(f"--- 🔍 REQUEST RECEIVED: GET /exam/{session_id} ---")

    # חיפוש המבחן במסד הנתונים
    exam_record = db.query(db_models.DBExam).filter(db_models.DBExam.id == session_id).first()

    if not exam_record:
        print("  -> ❌ ERROR: Exam not found in DB.")
        raise HTTPException(status_code=404, detail="המבחן לא נמצא או שהקישור פג תוקף.")

    # בדיקה האם יש למבחן הזה שאלות (אולי הוא רק נוצר אבל ה-AI נכשל באמצע)
    if not exam_record.questions:
        print("  -> ❌ ERROR: Exam exists but has no questions.")
        raise HTTPException(status_code=404, detail="המבחן נמצא אך אין בו שאלות.")

    try:
        response_questions = []

        # המרה של השאלות ממסד הנתונים חזרה למבנה שגיבשנו (Schemas)
        for q in exam_record.questions:
            response_questions.append(schemas.Question(
                id=q.id,
                type=q.type,
                question_text=q.question_text,
                options=q.options,
                correct_answer=q.correct_answer,
                explanation=q.explanation,
                difficulty=q.difficulty
            ))

        print(f"  -> ✅ Retrieved exam successfully with {len(response_questions)} questions!")
        print("=" * 40 + "\n")

        # החזרת המבחן המלא
        return schemas.Exam(
            id=exam_record.id,
            title=exam_record.title,
            questions=response_questions,
            source_material_summary=exam_record.source_material_summary
        )

    except Exception as e:
        print("\n🔥 --- CRITICAL ERROR IN /exam/{session_id} --- 🔥")
        print(f"Error Message: {str(e)}")
        raise HTTPException(status_code=500, detail="שגיאה בשליפת המבחן ממסד הנתונים.")